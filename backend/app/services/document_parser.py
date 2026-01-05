import os
import json
from typing import Dict, Any, Optional
from app.core.config import settings

# Try to import document processing libraries
try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None


async def parse_document(file_path: str, extension: str) -> Dict[str, Any]:
    """
    Parse a document and extract content.
    Returns dict with title, content, keywords, and summary.
    """
    extension = extension.lower()

    if extension == ".docx":
        return await parse_docx(file_path)
    elif extension == ".pptx":
        return await parse_pptx(file_path)
    elif extension == ".pdf":
        return await parse_pdf(file_path)
    elif extension in [".xlsx", ".xls"]:
        return await parse_excel(file_path)
    elif extension in [".md", ".txt"]:
        return await parse_text(file_path)
    elif extension == ".srt":
        return await parse_srt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {extension}")


async def parse_docx(file_path: str) -> Dict[str, Any]:
    """Parse DOCX file."""
    if DocxDocument is None:
        raise ImportError("python-docx is not installed")

    doc = DocxDocument(file_path)
    paragraphs = []

    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text.strip())

    content = "\n\n".join(paragraphs)

    # Extract title from first heading or filename
    title = os.path.basename(file_path).replace(".docx", "")
    for para in doc.paragraphs:
        if para.style and "Heading" in para.style.name:
            title = para.text.strip()
            break

    return {
        "title": title,
        "content": content,
        "keywords": extract_keywords(content),
        "summary": generate_simple_summary(content),
    }


async def parse_pptx(file_path: str) -> Dict[str, Any]:
    """Parse PPTX file."""
    if Presentation is None:
        raise ImportError("python-pptx is not installed")

    prs = Presentation(file_path)
    slides_content = []

    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())

        if slide_texts:
            slides_content.append(f"## 슬라이드 {i}\n" + "\n".join(slide_texts))

    content = "\n\n".join(slides_content)

    # Extract title from first slide
    title = os.path.basename(file_path).replace(".pptx", "")
    if prs.slides:
        first_slide = prs.slides[0]
        for shape in first_slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                title = shape.text.strip()[:100]
                break

    return {
        "title": title,
        "content": content,
        "keywords": extract_keywords(content),
        "summary": generate_simple_summary(content),
    }


async def parse_pdf(file_path: str) -> Dict[str, Any]:
    """Parse PDF file."""
    if fitz is None:
        raise ImportError("PyMuPDF is not installed")

    doc = fitz.open(file_path)
    pages_content = []

    for i, page in enumerate(doc, 1):
        text = page.get_text().strip()
        if text:
            pages_content.append(f"## 페이지 {i}\n{text}")

    content = "\n\n".join(pages_content)

    title = os.path.basename(file_path).replace(".pdf", "")

    return {
        "title": title,
        "content": content,
        "keywords": extract_keywords(content),
        "summary": generate_simple_summary(content),
    }


async def parse_excel(file_path: str) -> Dict[str, Any]:
    """Parse Excel file."""
    if load_workbook is None:
        raise ImportError("openpyxl is not installed")

    wb = load_workbook(file_path, data_only=True)
    sheets_content = []

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        rows = []

        for row in sheet.iter_rows(values_only=True):
            row_values = [str(cell) if cell is not None else "" for cell in row]
            if any(row_values):
                rows.append(" | ".join(row_values))

        if rows:
            sheets_content.append(f"## {sheet_name}\n" + "\n".join(rows))

    content = "\n\n".join(sheets_content)
    title = os.path.basename(file_path).replace(".xlsx", "").replace(".xls", "")

    return {
        "title": title,
        "content": content,
        "keywords": extract_keywords(content),
        "summary": generate_simple_summary(content),
    }


async def parse_text(file_path: str) -> Dict[str, Any]:
    """Parse text/markdown file."""
    with open(file_path, "r", encoding="utf-8") as f:
        content = f.read()

    title = os.path.basename(file_path).rsplit(".", 1)[0]

    # Try to extract title from markdown heading
    lines = content.split("\n")
    for line in lines:
        if line.startswith("# "):
            title = line[2:].strip()
            break

    return {
        "title": title,
        "content": content,
        "keywords": extract_keywords(content),
        "summary": generate_simple_summary(content),
    }


async def parse_srt(file_path: str) -> Dict[str, Any]:
    """Parse SRT subtitle file."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
    except UnicodeDecodeError:
        with open(file_path, "r", encoding="cp949") as f:
            content = f.read()

    # Extract only text lines (skip numbers and timestamps)
    lines = content.split("\n")
    text_lines = []

    for line in lines:
        line = line.strip()
        if not line:
            continue
        if line.isdigit():
            continue
        if "-->" in line:
            continue
        text_lines.append(line)

    text_content = " ".join(text_lines)
    title = os.path.basename(file_path).replace(".srt", "")

    return {
        "title": title,
        "content": content,  # Keep original SRT format
        "keywords": extract_keywords(text_content),
        "summary": generate_simple_summary(text_content),
    }


def extract_keywords(content: str, max_keywords: int = 10) -> list:
    """Extract simple keywords from content."""
    # Simple keyword extraction (word frequency)
    import re
    from collections import Counter

    # Remove special characters and split
    words = re.findall(r"\b[가-힣a-zA-Z]{2,}\b", content)

    # Filter common stop words (Korean + English)
    stop_words = {
        "이", "그", "저", "것", "수", "등", "및", "the", "a", "an", "is", "are",
        "was", "were", "be", "been", "being", "have", "has", "had", "do", "does",
        "did", "will", "would", "could", "should", "may", "might", "must",
        "있는", "하는", "되는", "있다", "하다", "되다", "위해", "통해", "대한",
    }

    filtered_words = [w for w in words if w.lower() not in stop_words]
    word_counts = Counter(filtered_words)

    return [word for word, _ in word_counts.most_common(max_keywords)]


def generate_simple_summary(content: str, max_length: int = 500) -> str:
    """Generate a simple summary (first few sentences)."""
    # Split into sentences
    sentences = content.replace("\n", " ").split(".")
    summary_parts = []
    current_length = 0

    for sentence in sentences:
        sentence = sentence.strip()
        if not sentence:
            continue
        if current_length + len(sentence) > max_length:
            break
        summary_parts.append(sentence)
        current_length += len(sentence)

    return ". ".join(summary_parts) + "." if summary_parts else content[:max_length]
